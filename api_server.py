import os
import platform
import tempfile
import zipfile
from typing import List, Optional

from deep_translator import GoogleTranslator
from docx import Document
from fastapi import FastAPI, File, Form, UploadFile
from fastapi.middleware.cors import CORSMiddleware
from fastapi.responses import FileResponse, JSONResponse
from openpyxl import load_workbook
from pdf2docx import Converter
from PIL import Image
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

app = FastAPI(title="Pro File Tool Master API", version="2.0.0")
app.add_middleware(CORSMiddleware, allow_origins=["*"], allow_credentials=True, allow_methods=["*"], allow_headers=["*"])
IS_WINDOWS = platform.system() == "Windows"


def save_upload(upload: UploadFile) -> str:
    ext = os.path.splitext(upload.filename or "upload.bin")[1]
    fd, path = tempfile.mkstemp(suffix=ext)
    os.close(fd)
    with open(path, "wb") as f:
        f.write(upload.file.read())
    return path


def fmt(n: int) -> str:
    if n < 1024: return f"{n} B"
    if n < 1024**2: return f"{n/1024:.1f} KB"
    return f"{n/1024**2:.2f} MB"


def mb(v) -> int:
    return int(float(v) * 1024 * 1024)


def zip_paths(paths: List[str], name="result") -> str:
    fd, out = tempfile.mkstemp(suffix=".zip"); os.close(fd)
    with zipfile.ZipFile(out, "w", zipfile.ZIP_DEFLATED, compresslevel=9) as z:
        for i, p in enumerate(paths, 1):
            z.write(p, arcname=f"{name}_{i:03d}{os.path.splitext(p)[1]}")
    return out


def txt_file(text: str, suffix=".txt") -> str:
    fd, out = tempfile.mkstemp(suffix=suffix); os.close(fd)
    with open(out, "w", encoding="utf-8") as f: f.write(text)
    return out


def translate_text(text, lang):
    if not text or not isinstance(text, str) or len(text.strip()) < 2 or text.isnumeric(): return text
    try: return GoogleTranslator(source="auto", target=lang).translate(text)
    except Exception: return text


def translate_docx(path, lang):
    doc = Document(path)
    for p in doc.paragraphs:
        if p.text.strip(): p.text = translate_text(p.text, lang)
    for t in doc.tables:
        for r in t.rows:
            for c in r.cells:
                for p in c.paragraphs:
                    if p.text.strip(): p.text = translate_text(p.text, lang)
    out = path.replace(".docx", f"_{lang}.docx"); doc.save(out); return out


def translate_xlsx(path, lang):
    wb = load_workbook(path)
    for sh in wb.worksheets:
        for row in sh.iter_rows():
            for cell in row:
                if cell.value and isinstance(cell.value, str) and not cell.value.startswith("="):
                    cell.value = translate_text(cell.value, lang)
    out = path.replace(".xlsx", f"_{lang}.xlsx"); wb.save(out); return out


def translate_pptx(path, lang):
    prs = Presentation(path)
    def shape(s):
        if s.shape_type == MSO_SHAPE_TYPE.GROUP:
            for ch in s.shapes: shape(ch)
        if getattr(s, "has_text_frame", False):
            for p in s.text_frame.paragraphs:
                if p.text.strip(): p.text = translate_text(p.text, lang)
        if getattr(s, "has_table", False):
            for r in s.table.rows:
                for c in r.cells:
                    if c.text: c.text = translate_text(c.text, lang)
    for sl in prs.slides:
        for s in sl.shapes: shape(s)
    out = path.replace(".pptx", f"_{lang}.pptx"); prs.save(out); return out


def images_to_pdf(paths):
    imgs = [Image.open(p).convert("RGB") for p in paths]
    fd, out = tempfile.mkstemp(suffix=".pdf"); os.close(fd)
    imgs[0].save(out, save_all=True, append_images=imgs[1:]); return out


def pdf_to_docx(path):
    out = path.replace(".pdf", ".docx")
    cv = Converter(path); cv.convert(out); cv.close(); return out


def docx_to_pdf(path):
    if not IS_WINDOWS: raise RuntimeError("DOCX → PDF cần Windows + Microsoft Word hoặc backend có LibreOffice.")
    from docx2pdf import convert
    out = os.path.abspath(path).replace(".docx", ".pdf"); convert(os.path.abspath(path), out); return out


def pptx_to_pdf(path):
    if not IS_WINDOWS: raise RuntimeError("PPTX → PDF cần Windows + Microsoft PowerPoint hoặc backend có LibreOffice.")
    import comtypes.client, pythoncom
    pythoncom.CoInitialize(); out = os.path.abspath(path).replace(".pptx", ".pdf")
    appp = comtypes.client.CreateObject("Powerpoint.Application", dynamic=True); appp.Visible = 1
    deck = appp.Presentations.Open(os.path.abspath(path)); deck.SaveAs(out, 32); deck.Close(); appp.Quit(); return out


def compress_image(path, quality=65, max_width=1800, output_format=None):
    img = Image.open(path); fmt0 = (output_format or img.format or "JPEG").upper().replace("JPG", "JPEG")
    if img.mode in ("RGBA", "P") and fmt0 == "JPEG": img = img.convert("RGB")
    w, h = img.size
    if w > int(max_width): img = img.resize((int(max_width), int(h * int(max_width) / w)), Image.Resampling.LANCZOS)
    ext = ".jpg" if fmt0 == "JPEG" else f".{fmt0.lower()}"; fd, out = tempfile.mkstemp(suffix=ext); os.close(fd)
    kw = {"optimize": True}
    if fmt0 in ("JPEG", "WEBP"): kw["quality"] = int(quality)
    if fmt0 == "PNG": kw["compress_level"] = 9
    img.save(out, format=fmt0, **kw); return out


def office_recompress(path):
    root, ext = os.path.splitext(path); out = f"{root}_compressed{ext}"
    with zipfile.ZipFile(path, "r") as zin, zipfile.ZipFile(out, "w", zipfile.ZIP_DEFLATED, compresslevel=9) as zout:
        for it in zin.infolist(): zout.writestr(it, zin.read(it.filename))
    return out


def raster_pdf(path, dpi=120, quality=60):
    import fitz
    root, _ = os.path.splitext(path); out = f"{root}_{dpi}dpi.pdf"
    doc = fitz.open(path); nd = fitz.open(); mat = fitz.Matrix(dpi/72, dpi/72)
    for pg in doc:
        pix = pg.get_pixmap(matrix=mat, alpha=False)
        nd.new_page(width=pg.rect.width, height=pg.rect.height).insert_image(pg.rect, stream=pix.tobytes("jpeg", jpg_quality=quality))
    nd.save(out, garbage=4, deflate=True, clean=True); nd.close(); doc.close(); return out


def compress_pdf_target(path, target_mb=100, keep_text=False):
    import fitz
    root, _ = os.path.splitext(path); target = mb(target_mb); cands=[]
    doc = fitz.open(path); opt = f"{root}_optimized.pdf"; doc.save(opt, garbage=4, deflate=True, clean=True); doc.close(); cands.append(opt)
    if keep_text and os.path.getsize(opt) <= target: return opt
    for dpi, q in [(180,80),(150,70),(120,60),(100,50),(90,40),(72,35),(60,30)]:
        p = raster_pdf(path, dpi, q); cands.append(p)
        if os.path.getsize(p) <= target: break
    under=[p for p in cands if os.path.getsize(p)<=target]
    return max(under, key=os.path.getsize) if under else min(cands, key=os.path.getsize)


def split_pdf(path, max_part_mb=100):
    import fitz
    lim=mb(max_part_mb); doc=fitz.open(path); parts=[]; cur=fitz.open(); idx=1
    def save(d, i):
        if len(d)==0: return None
        out=os.path.join(tempfile.gettempdir(), f"part_{i:03d}.pdf"); d.save(out, garbage=4, deflate=True, clean=True); return out
    for pi in range(len(doc)):
        test=fitz.open(); test.insert_pdf(cur); test.insert_pdf(doc, from_page=pi, to_page=pi)
        fd,tmp=tempfile.mkstemp(suffix=".pdf"); os.close(fd); test.save(tmp, garbage=4, deflate=True, clean=True); sz=os.path.getsize(tmp); test.close(); os.remove(tmp)
        if len(cur)>0 and sz>lim:
            parts.append(save(cur, idx)); cur.close(); cur=fitz.open(); idx+=1
        cur.insert_pdf(doc, from_page=pi, to_page=pi)
    p=save(cur, idx)
    if p: parts.append(p)
    cur.close(); doc.close(); return parts


def merge_pdfs(paths):
    import fitz
    outdoc=fitz.open()
    for p in paths: outdoc.insert_pdf(fitz.open(p))
    fd,out=tempfile.mkstemp(suffix=".pdf"); os.close(fd); outdoc.save(out); outdoc.close(); return out


def parse_pages(s):
    nums=[]
    for part in (s or "").replace(" ", "").split(','):
        if not part: continue
        if '-' in part:
            a,b=map(int,part.split('-',1)); nums += list(range(a-1,b))
        else: nums.append(int(part)-1)
    return sorted(set([n for n in nums if n>=0]))


def extract_pages(path, page_range):
    import fitz
    doc=fitz.open(path); nd=fitz.open()
    for i in parse_pages(page_range):
        if i < len(doc): nd.insert_pdf(doc, from_page=i, to_page=i)
    fd,out=tempfile.mkstemp(suffix=".pdf"); os.close(fd); nd.save(out); nd.close(); doc.close(); return out


def delete_pages(path, page_range):
    import fitz
    doc=fitz.open(path)
    for i in sorted(parse_pages(page_range), reverse=True):
        if i < len(doc): doc.delete_page(i)
    fd,out=tempfile.mkstemp(suffix=".pdf"); os.close(fd); doc.save(out); doc.close(); return out


def rotate_pdf(path, deg=90):
    import fitz
    doc=fitz.open(path)
    for p in doc: p.set_rotation((p.rotation + int(deg)) % 360)
    fd,out=tempfile.mkstemp(suffix=".pdf"); os.close(fd); doc.save(out); doc.close(); return out


def pdf_to_images(path, dpi=150):
    import fitz
    doc=fitz.open(path); outs=[]; mat=fitz.Matrix(dpi/72,dpi/72)
    for i,p in enumerate(doc,1):
        pix=p.get_pixmap(matrix=mat, alpha=False); fd,out=tempfile.mkstemp(suffix=f"_{i:03d}.jpg"); os.close(fd); pix.save(out); outs.append(out)
    doc.close(); return zip_paths(outs, "page")


def extract_text_pdf(path):
    import fitz
    doc=fitz.open(path); text="\n\n".join([p.get_text() for p in doc]); doc.close(); return txt_file(text or "No text layer found.")


def pdf_metadata(path):
    import fitz, json
    doc=fitz.open(path); data={"pages": len(doc), "metadata": doc.metadata}; doc.close(); return txt_file(json.dumps(data, ensure_ascii=False, indent=2))


def split_binary(path, part_mb=100):
    parts=[]; size=mb(part_mb); base=os.path.splitext(os.path.basename(path))[0]
    with open(path,"rb") as f:
        i=1
        while True:
            chunk=f.read(size)
            if not chunk: break
            out=os.path.join(tempfile.gettempdir(), f"{base}.part{i:03d}")
            with open(out,"wb") as w: w.write(chunk)
            parts.append(out); i+=1
    return zip_paths(parts, base)


def headers(inp, out, action, plan=""):
    ins=os.path.getsize(inp) if inp and os.path.exists(inp) else 0; outs=os.path.getsize(out)
    return {"X-Tool-Action": action, "X-Master-Plan": plan, "X-Input-Size-Text": fmt(ins), "X-Output-Size-Text": fmt(outs), "X-Saved-Percent": f"{max((ins-outs)/ins*100,0):.1f}" if ins else "0"}


@app.get("/health")
def health(): return {"ok": True, "service": "pro-file-tool-master-api", "tools": 25}


@app.post("/tool")
def tool(action: str = Form(...), goal: str = Form("ai_read_large_file"), file: UploadFile = File(None), files: List[UploadFile] = File(None), target_lang: str = Form("vi"), target_mb: float = Form(100), max_part_mb: float = Form(100), page_range: str = Form("1-1"), rotate_degrees: int = Form(90), output_format: str = Form("JPEG"), max_width: int = Form(1800)):
    try:
        paths=[save_upload(f) for f in (files or [])]
        main=save_upload(file) if file else (paths[0] if paths else None)
        ext=os.path.splitext(file.filename if file else (files[0].filename if files else "result"))[1].lower()
        base=os.path.splitext(file.filename if file else (files[0].filename if files else "result"))[0]
        out=None; name=f"{base}_result"; plan=action
        if action=="master_tool":
            if goal=="ai_read_large_file":
                plan="compress_pdf_to_target -> split_pdf_by_size -> zip" if ext==".pdf" else "compress/split by type"
                out=zip_paths(split_pdf(compress_pdf_target(main,target_mb,False),max_part_mb), base) if ext==".pdf" else split_binary(main,max_part_mb); name=f"{base}_master_ai.zip"
            elif goal=="smallest_possible": out=compress_pdf_target(main,target_mb,False) if ext==".pdf" else compress_image(main,40,1200); name=f"{base}_small{os.path.splitext(out)[1]}"
            elif goal=="keep_quality": out=compress_pdf_target(main,target_mb,True) if ext==".pdf" else compress_image(main,85,2400); name=f"{base}_quality{os.path.splitext(out)[1]}"
            elif goal=="translate_document": action={".docx":"translate_docx",".xlsx":"translate_xlsx",".pptx":"translate_pptx"}.get(ext,"translate_docx")
            elif goal=="convert_for_sharing": action="images_to_pdf" if ext in [".jpg",".jpeg",".png",".webp"] else "pdf_to_docx"
            elif goal=="extract_for_ai": out=zip_paths([extract_text_pdf(main)] + (split_pdf(compress_pdf_target(main,target_mb,False),max_part_mb) if ext==".pdf" else []), base); name=f"{base}_ai_extract.zip"
        if out is None:
            if action=="translate_docx": out=translate_docx(main,target_lang); name=f"{base}_{target_lang}.docx"
            elif action=="translate_xlsx": out=translate_xlsx(main,target_lang); name=f"{base}_{target_lang}.xlsx"
            elif action=="translate_pptx": out=translate_pptx(main,target_lang); name=f"{base}_{target_lang}.pptx"
            elif action=="images_to_pdf": out=images_to_pdf(paths or [main]); name=f"{base}_images.pdf"
            elif action=="pdf_to_docx": out=pdf_to_docx(main); name=f"{base}.docx"
            elif action=="docx_to_pdf": out=docx_to_pdf(main); name=f"{base}.pdf"
            elif action=="pptx_to_pdf": out=pptx_to_pdf(main); name=f"{base}.pdf"
            elif action=="compress": out=compress_pdf_target(main,target_mb,False) if ext==".pdf" else (office_recompress(main) if ext in [".docx",".xlsx",".pptx"] else compress_image(main)); name=f"{base}_compressed{os.path.splitext(out)[1]}"
            elif action=="compress_split_pdf": out=zip_paths(split_pdf(compress_pdf_target(main,target_mb,False),max_part_mb), base); name=f"{base}_ai_parts.zip"
            elif action=="split_pdf": out=zip_paths(split_pdf(main,max_part_mb), base); name=f"{base}_parts.zip"
            elif action=="merge_pdfs": out=merge_pdfs(paths); name=f"{base}_merged.pdf"
            elif action=="extract_pages": out=extract_pages(main,page_range); name=f"{base}_pages.pdf"
            elif action=="delete_pages": out=delete_pages(main,page_range); name=f"{base}_deleted.pdf"
            elif action=="rotate_pdf": out=rotate_pdf(main,rotate_degrees); name=f"{base}_rotated.pdf"
            elif action=="pdf_to_images": out=pdf_to_images(main); name=f"{base}_images.zip"
            elif action=="extract_text_pdf": out=extract_text_pdf(main); name=f"{base}.txt"
            elif action=="pdf_metadata": out=pdf_metadata(main); name=f"{base}_metadata.txt"
            elif action=="convert_image": out=compress_image(main,85,max_width,output_format); name=f"{base}.{output_format.lower().replace('jpeg','jpg')}"
            elif action=="resize_image": out=compress_image(main,85,max_width); name=f"{base}_resized{os.path.splitext(out)[1]}"
            elif action=="compress_image": out=compress_image(main,60,max_width); name=f"{base}_compressed{os.path.splitext(out)[1]}"
            elif action=="zip_files": out=zip_paths(paths,base); name=f"{base}_files.zip"
            elif action=="split_binary": out=split_binary(main,max_part_mb); name=f"{base}_binary_parts.zip"
            else: return JSONResponse({"error":"Unknown action"}, status_code=400)
        return FileResponse(out, filename=name, headers=headers(main,out,action,plan))
    except Exception as e:
        return JSONResponse({"error": str(e)}, status_code=500)


@app.post("/process")
def process(file: UploadFile = File(...), target_mb: float = Form(100), max_part_mb: float = Form(100), split_for_ai: bool = Form(True), pdf_mode: str = Form("deep")):
    return tool("compress_split_pdf" if split_for_ai else "compress", "ai_read_large_file", file, None, "vi", target_mb, max_part_mb)
