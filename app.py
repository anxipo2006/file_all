import streamlit as st
import os
import tempfile
import platform
import zipfile
from PIL import Image
from pdf2docx import Converter
from docx import Document
from openpyxl import load_workbook
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from deep_translator import GoogleTranslator

# --- CẤU HÌNH TRANG ---
st.set_page_config(
    page_title="Pro File Tool | Đa Năng", 
    page_icon="⚡", 
    layout="wide",
    initial_sidebar_state="expanded"
)

# --- CSS GIAO DIỆN (MODERN DARK MODE) ---
st.markdown("""
<style>
    /* Tổng quan */
    .main-header { 
        text-align: center; 
        background: -webkit-linear-gradient(45deg, #00e676, #00b0ff);
        -webkit-background-clip: text;
        -webkit-text-fill-color: transparent;
        font-size: 3rem; 
        font-weight: 900; 
        margin-bottom: 5px; 
    }
    .sub-header { text-align: center; color: #888; margin-bottom: 30px; font-size: 1.1rem; }
    
    /* Card Style */
    .css-card {
        background: rgba(38, 39, 48, 0.7);
        backdrop-filter: blur(10px);
        border-radius: 15px; 
        padding: 20px;
        box-shadow: 0 8px 32px 0 rgba(0, 0, 0, 0.3);
        border: 1px solid rgba(255, 255, 255, 0.1);
        text-align: center; 
        margin-bottom: 15px; 
        transition: all 0.3s ease;
    }
    .css-card:hover { 
        transform: translateY(-5px); 
        border-color: #00e676; 
        box-shadow: 0 8px 32px 0 rgba(0, 230, 118, 0.2);
    }
    
    /* Headers & Text */
    .card-header { color: #fff; font-size: 1.2rem; font-weight: 700; border-bottom: 1px solid #444; padding-bottom: 12px; margin-bottom: 15px;}
    
    /* Buttons */
    .stButton>button { width: 100%; border-radius: 8px; font-weight: bold; transition: 0.3s; }
    .stDownloadButton>button { background-color: #00e676; color: #000; }
    .stDownloadButton>button:hover { background-color: #00c853; color: #000; }
</style>
""", unsafe_allow_html=True)

# --- HÀM HỖ TRỢ HỆ THỐNG ---
def save_uploaded_file(uploaded_file):
    """Lưu file tạm thời và trả về đường dẫn"""
    try:
        ext = os.path.splitext(uploaded_file.name)[1]
        with tempfile.NamedTemporaryFile(delete=False, suffix=ext) as tmp:
            tmp.write(uploaded_file.getbuffer())
            return tmp.name
    except Exception as e:
        st.error(f"Lỗi khi lưu file: {e}")
        return None

def cleanup_files(*filepaths):
    """Xóa file tạm để giải phóng bộ nhớ server"""
    for path in filepaths:
        if path and os.path.exists(path):
            try:
                os.remove(path)
            except:
                pass

IS_WINDOWS = platform.system() == "Windows"

# ==========================================
# PHẦN 1: LOGIC DỊCH THUẬT
# ==========================================

def standard_translate(text, target_lang):
    """Dịch an toàn, bỏ qua chuỗi rỗng hoặc số"""
    if not text or not isinstance(text, str) or len(text.strip()) < 2 or text.isnumeric():
        return text
    try:
        translator = GoogleTranslator(source='auto', target=target_lang)
        return translator.translate(text)
    except:
        return text

def process_word_std(docx_path, target_lang, status_placeholder):
    doc = Document(docx_path)
    total = len(doc.paragraphs) + sum([len(t.rows) for t in doc.tables])
    if total == 0: total = 1
    processed = 0
    
    def apply_trans(para):
        if para.text.strip():
            style = para.runs[0] if para.runs else None
            para.text = standard_translate(para.text, target_lang)
            if style and para.runs:
                r = para.runs[0]
                r.font.name = style.font.name
                r.font.size = style.font.size
                r.font.bold = style.font.bold
                r.font.italic = style.font.italic
                if hasattr(style.font, 'color') and hasattr(style.font.color, 'rgb') and style.font.color.rgb:
                    try: r.font.color.rgb = style.font.color.rgb
                    except: pass

    for i, para in enumerate(doc.paragraphs):
        apply_trans(para)
        processed += 1
        if i % 10 == 0: status_placeholder.progress(min(processed/total, 0.9), text=f"Đang dịch nội dung...")

    for table in doc.tables:
        for row in table.rows:
            for cell in row.cells:
                for para in cell.paragraphs: 
                    apply_trans(para)
                    processed += 1
    
    status_placeholder.progress(1.0, text="Hoàn tất xử lý Word!")
    out = docx_path.replace(".docx", f"_{target_lang.upper()}.docx")
    doc.save(out)
    return out

def process_excel_std(xlsx_path, target_lang, status_placeholder):
    wb = load_workbook(xlsx_path)
    total_sheets = len(wb.worksheets)
    
    for i, sheet in enumerate(wb.worksheets):
        status_placeholder.progress((i+1)/total_sheets, text=f"Đang dịch Sheet: {sheet.title}")
        for row in sheet.iter_rows():
            for cell in row:
                if cell.value and isinstance(cell.value, str) and not cell.value.startswith("="):
                    cell.value = standard_translate(cell.value, target_lang)
    
    status_placeholder.progress(1.0, text="Hoàn tất xử lý Excel!")
    out = xlsx_path.replace(".xlsx", f"_{target_lang.upper()}.xlsx")
    wb.save(out)
    return out

def process_ppt_std(ppt_path, target_lang, status_placeholder):
    prs = Presentation(ppt_path)
    total_slides = len(prs.slides)
    if total_slides == 0: total_slides = 1
    
    def process_shape(shape):
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            for child in shape.shapes: process_shape(child)
        
        if shape.has_text_frame:
            for p in shape.text_frame.paragraphs:
                if p.text.strip():
                    # Lưu format cũ
                    font_info = {}
                    if p.runs:
                        r = p.runs[0]
                        font_info = {'name': r.font.name, 'size': r.font.size, 'bold': r.font.bold, 'italic': r.font.italic}
                        try: font_info['color'] = r.font.color.rgb
                        except: pass
                    
                    trans = standard_translate(p.text, target_lang)
                    p.clear()
                    run = p.add_run()
                    run.text = trans
                    
                    # Áp dụng lại format
                    if font_info.get('name'): run.font.name = font_info['name']
                    if font_info.get('size'): run.font.size = font_info['size']
                    run.font.bold = font_info.get('bold')
                    run.font.italic = font_info.get('italic')
                    if font_info.get('color'): run.font.color.rgb = font_info['color']
        
        if shape.has_table:
            for row in shape.table.rows:
                for cell in row.cells:
                    if cell.text_frame:
                         for p in cell.text_frame.paragraphs:
                             if p.text.strip():
                                 p.text = standard_translate(p.text, target_lang)

    for i, slide in enumerate(prs.slides):
        status_placeholder.progress((i+1)/total_slides, text=f"Đang dịch Slide {i+1}/{total_slides}...")
        for shape in slide.shapes: process_shape(shape)
        
    status_placeholder.progress(1.0, text="Hoàn tất xử lý PPT!")
    out = ppt_path.replace(".pptx", f"_{target_lang.upper()}.pptx")
    prs.save(out)
    return out

# ==========================================
# PHẦN 2: LOGIC CHUYỂN ĐỔI (CONVERTER)
# ==========================================

def process_img_to_pdf(files):
    try:
        imgs = [Image.open(f).convert('RGB') for f in files]
        tmp_fd, tmp_path = tempfile.mkstemp(suffix=".pdf")
        os.close(tmp_fd)
        imgs[0].save(tmp_path, save_all=True, append_images=imgs[1:])
        return tmp_path
    except Exception as e:
        st.error(f"Lỗi ảnh: {e}")
        return None

def process_pdf_to_word(path):
    try:
        docx_path = path.replace(".pdf", ".docx")
        cv = Converter(path)
        cv.convert(docx_path)
        cv.close()
        return docx_path
    except Exception as e:
        st.error(f"Lỗi convert PDF: {e}")
        return None

def process_word_to_pdf(path):
    if IS_WINDOWS:
        try:
            from docx2pdf import convert
            abs_path = os.path.abspath(path)
            pdf_path = abs_path.replace(".docx", ".pdf")
            convert(abs_path, pdf_path)
            return pdf_path
        except Exception as e:
            st.error(f"Lỗi MS Word Engine: {e}")
            return None
    return None

def process_ppt_to_pdf(path):
    if IS_WINDOWS:
        try:
            import comtypes.client
            import pythoncom
            pythoncom.CoInitialize() # Yêu cầu cho thread Streamlit
            abs_path = os.path.abspath(path)
            pdf_path = abs_path.replace(".pptx", ".pdf").replace(".ppt", ".pdf")
            ppt_app = comtypes.client.CreateObject("Powerpoint.Application", dynamic=True)
            ppt_app.Visible = 1
            deck = ppt_app.Presentations.Open(abs_path)
            deck.SaveAs(pdf_path, 32)
            deck.Close()
            ppt_app.Quit()
            return pdf_path
        except Exception as e:
            st.error(f"Lỗi MS PowerPoint Engine: {e}")
            return None
    return None

# ==========================================
# PHẦN 3: LOGIC HẠ DUNG LƯỢNG FILE
# ==========================================

def format_file_size(size_bytes):
    """Định dạng dung lượng file dễ đọc"""
    if size_bytes < 1024:
        return f"{size_bytes} B"
    if size_bytes < 1024 * 1024:
        return f"{size_bytes / 1024:.1f} KB"
    return f"{size_bytes / (1024 * 1024):.2f} MB"

def get_compression_profile(level):
    """Thông số nén theo mức người dùng chọn"""
    profiles = {
        "Nhẹ - giữ chất lượng cao": {"quality": 85, "max_width": 2400, "pdf_dpi": 160},
        "Cân bằng - khuyên dùng": {"quality": 65, "max_width": 1800, "pdf_dpi": 120},
        "Tối đa - dung lượng nhỏ": {"quality": 40, "max_width": 1200, "pdf_dpi": 90},
    }
    return profiles.get(level, profiles["Cân bằng - khuyên dùng"])

def process_image_compress(uploaded_file, quality=70, max_width=1800):
    """Nén ảnh, tự resize nếu ảnh quá rộng"""
    try:
        img = Image.open(uploaded_file)
        img_format = (img.format or "JPEG").upper()
        if img.mode in ("RGBA", "P") and img_format in ("JPEG", "JPG"):
            img = img.convert("RGB")

        width, height = img.size
        if width > max_width:
            new_height = int(height * (max_width / width))
            img = img.resize((max_width, new_height), Image.Resampling.LANCZOS)

        ext = ".jpg" if img_format in ("JPEG", "JPG") else f".{img_format.lower()}"
        tmp_fd, tmp_path = tempfile.mkstemp(suffix=ext)
        os.close(tmp_fd)

        save_kwargs = {"optimize": True}
        if img_format in ("JPEG", "JPG", "WEBP"):
            save_kwargs["quality"] = quality
        if img_format == "PNG":
            save_kwargs["compress_level"] = 9

        img.save(tmp_path, format=img_format, **save_kwargs)
        return tmp_path
    except Exception as e:
        st.error(f"Lỗi nén ảnh: {e}")
        return None

def process_pdf_compress(path, dpi=120, image_quality=65):
    """Nén PDF mạnh: tối ưu cấu trúc, nếu chưa giảm thì rasterize từng trang"""
    try:
        import fitz
        root, _ = os.path.splitext(path)
        optimized_out = f"{root}_optimized.pdf"
        raster_out = f"{root}_compressed.pdf"

        doc = fitz.open(path)
        doc.save(optimized_out, garbage=4, deflate=True, clean=True)
        original_size = os.path.getsize(path)
        optimized_size = os.path.getsize(optimized_out)

        # Nếu tối ưu cấu trúc đã giảm tốt, giữ nguyên text/vector để chất lượng cao hơn.
        if optimized_size < original_size * 0.85:
            doc.close()
            return optimized_out

        # Với PDF scan/ảnh lớn: dựng lại từng trang thành JPEG để giảm mạnh dung lượng.
        new_doc = fitz.open()
        zoom = dpi / 72
        matrix = fitz.Matrix(zoom, zoom)
        total_pages = max(len(doc), 1)
        progress = st.progress(0, text="Đang nén sâu PDF...")

        for index, page in enumerate(doc):
            pix = page.get_pixmap(matrix=matrix, alpha=False)
            img_bytes = pix.tobytes("jpeg", jpg_quality=image_quality)
            new_page = new_doc.new_page(width=page.rect.width, height=page.rect.height)
            new_page.insert_image(page.rect, stream=img_bytes)
            progress.progress((index + 1) / total_pages, text=f"Đang nén sâu PDF {index + 1}/{total_pages}...")

        new_doc.save(raster_out, garbage=4, deflate=True, clean=True)
        new_doc.close()
        doc.close()
        progress.empty()

        if os.path.exists(optimized_out):
            os.remove(optimized_out)
        return raster_out
    except Exception as e:
        st.error(f"Lỗi nén PDF: {e}")
        return None

def process_office_recompress(path):
    """Nén lại file Office OpenXML (.docx/.xlsx/.pptx) bằng ZIP deflate"""
    try:
        root, ext = os.path.splitext(path)
        out = f"{root}_compressed{ext}"
        with zipfile.ZipFile(path, "r") as zin:
            with zipfile.ZipFile(out, "w", compression=zipfile.ZIP_DEFLATED, compresslevel=9) as zout:
                for item in zin.infolist():
                    zout.writestr(item, zin.read(item.filename))
        return out
    except Exception as e:
        st.error(f"Lỗi nén Office: {e}")
        return None

def get_compression_stats(input_size, output_path):
    """Tính dung lượng sau nén và phần trăm giảm"""
    output_size = os.path.getsize(output_path)
    saved = input_size - output_size
    ratio = (saved / input_size * 100) if input_size else 0
    return output_size, saved, ratio

def mb_to_bytes(mb):
    """Đổi MB sang bytes"""
    return int(float(mb) * 1024 * 1024)

def make_zip(filepaths, output_name="ai_parts.zip"):
    """Gom nhiều file thành ZIP để tải một lần"""
    tmp_fd, zip_path = tempfile.mkstemp(suffix=".zip")
    os.close(tmp_fd)
    with zipfile.ZipFile(zip_path, "w", compression=zipfile.ZIP_DEFLATED, compresslevel=9) as zf:
        for path in filepaths:
            zf.write(path, arcname=os.path.basename(path))
    return zip_path

def rasterize_pdf(path, dpi=120, image_quality=65, suffix="compressed"):
    """Dựng lại PDF thành ảnh JPEG theo từng trang để giảm mạnh dung lượng"""
    try:
        import fitz
        root, _ = os.path.splitext(path)
        out = f"{root}_{suffix}.pdf"
        doc = fitz.open(path)
        new_doc = fitz.open()
        zoom = dpi / 72
        matrix = fitz.Matrix(zoom, zoom)
        total_pages = max(len(doc), 1)
        progress = st.progress(0, text="Đang nén sâu PDF...")

        for index, page in enumerate(doc):
            pix = page.get_pixmap(matrix=matrix, alpha=False)
            img_bytes = pix.tobytes("jpeg", jpg_quality=image_quality)
            new_page = new_doc.new_page(width=page.rect.width, height=page.rect.height)
            new_page.insert_image(page.rect, stream=img_bytes)
            progress.progress((index + 1) / total_pages, text=f"Đang nén sâu PDF {index + 1}/{total_pages}...")

        new_doc.save(out, garbage=4, deflate=True, clean=True)
        new_doc.close()
        doc.close()
        progress.empty()
        return out
    except Exception as e:
        st.error(f"Lỗi rasterize PDF: {e}")
        return None

def compress_pdf_to_target(path, target_mb=None, mode="Giữ text nếu có thể"):
    """Nén PDF về gần dung lượng mục tiêu nhất có thể"""
    try:
        import fitz
        root, _ = os.path.splitext(path)
        target_bytes = mb_to_bytes(target_mb) if target_mb else None
        candidates = []

        doc = fitz.open(path)
        optimized_out = f"{root}_optimized.pdf"
        doc.save(optimized_out, garbage=4, deflate=True, clean=True)
        doc.close()
        candidates.append(optimized_out)

        if mode == "Giữ text nếu có thể" and (not target_bytes or os.path.getsize(optimized_out) <= target_bytes):
            return optimized_out

        presets = [
            (180, 80), (150, 70), (120, 60),
            (100, 50), (90, 40), (72, 35), (60, 30)
        ]
        for dpi, quality in presets:
            out = rasterize_pdf(path, dpi=dpi, image_quality=quality, suffix=f"{dpi}dpi_q{quality}")
            if out:
                candidates.append(out)
                if target_bytes and os.path.getsize(out) <= target_bytes:
                    break

        best = min(candidates, key=lambda p: os.path.getsize(p))
        if target_bytes:
            under_target = [p for p in candidates if os.path.getsize(p) <= target_bytes]
            if under_target:
                best = max(under_target, key=lambda p: os.path.getsize(p))

        for candidate in candidates:
            if candidate != best and os.path.exists(candidate):
                os.remove(candidate)
        return best
    except Exception as e:
        st.error(f"Lỗi nén PDF theo mục tiêu: {e}")
        return None

def split_pdf_by_size(path, max_part_mb=100):
    """Tách PDF theo trang để mỗi part gần dưới giới hạn MB"""
    try:
        import fitz
        max_bytes = mb_to_bytes(max_part_mb)
        doc = fitz.open(path)
        part_paths = []
        current = fitz.open()
        part_index = 1
        progress = st.progress(0, text="Đang tách PDF cho AI...")

        def save_part(part_doc, index):
            if len(part_doc) == 0:
                return None
            out = os.path.join(tempfile.gettempdir(), f"ai_part_{index:03d}.pdf")
            part_doc.save(out, garbage=4, deflate=True, clean=True)
            return out

        for page_index in range(len(doc)):
            test_doc = fitz.open()
            test_doc.insert_pdf(current)
            test_doc.insert_pdf(doc, from_page=page_index, to_page=page_index)
            tmp_fd, tmp_path = tempfile.mkstemp(suffix=".pdf")
            os.close(tmp_fd)
            test_doc.save(tmp_path, garbage=4, deflate=True, clean=True)
            test_size = os.path.getsize(tmp_path)
            test_doc.close()
            os.remove(tmp_path)

            if len(current) > 0 and test_size > max_bytes:
                saved = save_part(current, part_index)
                if saved:
                    part_paths.append(saved)
                current.close()
                current = fitz.open()
                part_index += 1

            current.insert_pdf(doc, from_page=page_index, to_page=page_index)
            progress.progress((page_index + 1) / max(len(doc), 1), text=f"Đang tách trang {page_index + 1}/{len(doc)}...")

        saved = save_part(current, part_index)
        if saved:
            part_paths.append(saved)
        current.close()
        doc.close()
        progress.empty()
        return part_paths
    except Exception as e:
        st.error(f"Lỗi tách PDF: {e}")
        return []

def compress_and_split_pdf(path, target_total_mb, max_part_mb, mode):
    """Nén PDF theo target tổng rồi tách parts cho AI"""
    compressed = compress_pdf_to_target(path, target_total_mb, mode)
    if not compressed:
        return None, [], None
    parts = split_pdf_by_size(compressed, max_part_mb)
    zip_path = make_zip(parts, "ai_parts.zip") if parts else None
    return compressed, parts, zip_path

# ==========================================
# GIAO DIỆN CHÍNH
# ==========================================

st.markdown("<div class='main-header'>⚡ PRO FILE TOOL</div>", unsafe_allow_html=True)
st.markdown("<div class='sub-header'>Công cụ xử lý tài liệu mạnh mẽ: Dịch thuật giữ nguyên định dạng & Chuyển đổi siêu tốc</div>", unsafe_allow_html=True)

tab_trans, tab_convert, tab_compress = st.tabs(["🌐 DỊCH THUẬT (SMART)", "📂 CHUYỂN ĐỔI FILE", "🗜️ HẠ DUNG LƯỢNG"])

# === TAB 1: DỊCH THUẬT ===
with tab_trans:
    lang_col1, lang_col2 = st.columns([1, 2])
    with lang_col1:
        lang_dict = {
            "Tiếng Việt": "vi", "Tiếng Anh": "en", "Tiếng Nhật": "ja", 
            "Tiếng Hàn": "ko", "Tiếng Trung": "zh-CN", "Tiếng Pháp": "fr"
        }
        target_lang_name = st.selectbox("🌐 Chọn ngôn ngữ đích:", list(lang_dict.keys()))
        target_lang = lang_dict[target_lang_name]

    st.info("💡 Hệ thống sẽ tự động nhận diện ngôn ngữ gốc và dịch sang ngôn ngữ bạn chọn, bảo toàn tối đa định dạng file.")
    
    tc1, tc2, tc3 = st.columns(3)
    
    # --- Xử lý Word ---
    with tc1:
        st.markdown("<div class='css-card'><div class='card-header'>📘 Word (.docx)</div></div>", unsafe_allow_html=True)
        f_w = st.file_uploader("Kéo thả Word vào đây", type=['docx'], key="s_w")
        if f_w and st.button("Bắt đầu dịch Word", key="b_s_w"):
            with st.spinner("Đang chuẩn bị dữ liệu..."):
                p_in = save_uploaded_file(f_w)
                bar = st.empty()
                p_out = process_word_std(p_in, target_lang, bar)
                
            with open(p_out, "rb") as f: 
                st.download_button(f"⬇️ Tải file ({target_lang.upper()})", f, f"{target_lang.upper()}_{f_w.name}")
                st.toast("Dịch Word thành công!", icon="✅")
            cleanup_files(p_in, p_out) # Dọn rác

    # --- Xử lý Excel ---
    with tc2:
        st.markdown("<div class='css-card'><div class='card-header'>📗 Excel (.xlsx)</div></div>", unsafe_allow_html=True)
        f_e = st.file_uploader("Kéo thả Excel vào đây", type=['xlsx'], key="s_e")
        if f_e and st.button("Bắt đầu dịch Excel", key="b_s_e"):
            with st.spinner("Đang xử lý các Sheet..."):
                p_in = save_uploaded_file(f_e)
                bar = st.empty()
                p_out = process_excel_std(p_in, target_lang, bar)
                
            with open(p_out, "rb") as f: 
                st.download_button(f"⬇️ Tải file ({target_lang.upper()})", f, f"{target_lang.upper()}_{f_e.name}")
                st.toast("Dịch Excel thành công!", icon="✅")
            cleanup_files(p_in, p_out)

    # --- Xử lý PPT ---
    with tc3:
        st.markdown("<div class='css-card'><div class='card-header'>📙 PowerPoint (.pptx)</div></div>", unsafe_allow_html=True)
        f_p = st.file_uploader("Kéo thả PPT vào đây", type=['pptx'], key="s_p")
        if f_p and st.button("Bắt đầu dịch PPT", key="b_s_p"):
            with st.spinner("Đang bóc tách Slides..."):
                p_in = save_uploaded_file(f_p)
                bar = st.empty()
                p_out = process_ppt_std(p_in, target_lang, bar)
                
            if p_out:
                with open(p_out, "rb") as f: 
                    st.download_button(f"⬇️ Tải file ({target_lang.upper()})", f, f"{target_lang.upper()}_{f_p.name}")
                    st.toast("Dịch PPT thành công!", icon="✅")
                cleanup_files(p_in, p_out)

# === TAB 2: CONVERTER ===
with tab_convert:
    if not IS_WINDOWS:
        st.warning("⚠️ Chức năng chuyển đổi Word ➡ PDF và PPT ➡ PDF yêu cầu hệ điều hành Windows (Sử dụng MS Office Engine). Môi trường hiện tại không hỗ trợ chức năng này.")

    col1, col2 = st.columns(2, gap="large")
    
    with col1:
        st.markdown("<div class='css-card'><div class='card-header'>🖼️ Image ➡ PDF</div></div>", unsafe_allow_html=True)
        u_img = st.file_uploader("Upload Ảnh (Có thể chọn nhiều)", type=['png','jpg','jpeg'], accept_multiple_files=True, key="c_img")
        if st.button("Convert Img ➡ PDF", key="btn_c1", disabled=not bool(u_img)):
            with st.spinner("Đang gộp ảnh..."):
                res = process_img_to_pdf(u_img)
                if res:
                    with open(res, "rb") as f: st.download_button("⬇️ Tải PDF", f, "images_converted.pdf")
                    cleanup_files(res)

    with col2:
        st.markdown("<div class='css-card'><div class='card-header'>📝 Word ➡ PDF</div></div>", unsafe_allow_html=True)
        u_word = st.file_uploader("Upload Word (.docx)", type=['docx'], key="c_word", disabled=not IS_WINDOWS)
        if st.button("Convert Word ➡ PDF", key="btn_c2", disabled=(not IS_WINDOWS or not bool(u_word))):
            with st.spinner("Đang gọi MS Word Engine..."):
                path_in = save_uploaded_file(u_word)
                res = process_word_to_pdf(path_in)
                if res:
                    with open(res, "rb") as f: st.download_button("⬇️ Tải PDF", f, f"{os.path.splitext(u_word.name)[0]}.pdf")
                    cleanup_files(path_in, res)

    st.markdown("---")
    col3, col4 = st.columns(2, gap="large")
    
    with col3:
        st.markdown("<div class='css-card'><div class='card-header'>📄 PDF ➡ Word</div></div>", unsafe_allow_html=True)
        u_pdf = st.file_uploader("Upload PDF (.pdf)", type=['pdf'], key="c_pdf")
        if st.button("Convert PDF ➡ Word", key="btn_c3", disabled=not bool(u_pdf)):
            with st.spinner("Đang bóc tách văn bản..."):
                path_in = save_uploaded_file(u_pdf)
                res = process_pdf_to_word(path_in)
                if res:
                    with open(res, "rb") as f: st.download_button("⬇️ Tải Word", f, f"{os.path.splitext(u_pdf.name)[0]}.docx")
                    cleanup_files(path_in, res)

    with col4:
        st.markdown("<div class='css-card'><div class='card-header'>📊 PPT ➡ PDF</div></div>", unsafe_allow_html=True)
        u_ppt = st.file_uploader("Upload PPT (.pptx)", type=['pptx'], key="c_ppt", disabled=not IS_WINDOWS)
        if st.button("Convert PPT ➡ PDF", key="btn_c4", disabled=(not IS_WINDOWS or not bool(u_ppt))):
            with st.spinner("Đang gọi MS PowerPoint Engine..."):
                path_in = save_uploaded_file(u_ppt)
                res = process_ppt_to_pdf(path_in)
                if res:
                    with open(res, "rb") as f: st.download_button("⬇️ Tải PDF", f, f"{os.path.splitext(u_ppt.name)[0]}.pdf")
                    cleanup_files(path_in, res)

# === TAB 3: HẠ DUNG LƯỢNG FILE ===
with tab_compress:
    st.info("💡 Đã tăng giới hạn upload lên 2048MB. Nén PDF scan sẽ dùng chế độ nén sâu để giảm tối đa dung lượng.")

    cc1, cc2 = st.columns([1, 1], gap="large")
    with cc1:
        st.markdown("<div class='css-card'><div class='card-header'>🗜️ Nén file</div></div>", unsafe_allow_html=True)
        u_comp = st.file_uploader(
            "Upload file cần hạ dung lượng",
            type=['jpg', 'jpeg', 'png', 'webp', 'pdf', 'docx', 'xlsx', 'pptx'],
            key="compress_file"
        )
        level = st.select_slider(
            "Mức nén nhanh",
            options=["Nhẹ - giữ chất lượng cao", "Cân bằng - khuyên dùng", "Tối đa - dung lượng nhỏ"],
            value="Cân bằng - khuyên dùng"
        )
        target_total_mb = st.number_input("Dung lượng muốn hạ về (MB)", min_value=1, max_value=2048, value=100, step=10)
        split_for_ai = st.checkbox("Chia thành nhiều file nhỏ cho AI", value=True)
        max_part_mb = st.number_input("Giới hạn mỗi file nhỏ (MB)", min_value=1, max_value=500, value=100, step=10)
        pdf_mode = st.radio(
            "Chế độ PDF",
            ["Giữ text nếu có thể", "Nén sâu tối đa cho AI đọc"],
            index=1
        )

    with cc2:
        st.markdown("<div class='css-card'><div class='card-header'>📊 Kết quả</div></div>", unsafe_allow_html=True)
        if u_comp:
            input_size = len(u_comp.getbuffer())
            st.metric("Dung lượng gốc", format_file_size(input_size))
            st.caption(f"Mục tiêu tổng: {target_total_mb} MB • Mỗi part tối đa: {max_part_mb} MB")
        else:
            st.caption("Upload file để xem kết quả nén.")

    if st.button("🚀 Hạ dung lượng / chia file", key="btn_compress", disabled=not bool(u_comp)):
        with st.spinner("Đang tối ưu file..."):
            ext = os.path.splitext(u_comp.name)[1].lower()
            profile = get_compression_profile(level)
            p_in = None
            p_out = None
            part_paths = []
            zip_path = None
            cleanup_later = []

            if ext in ['.jpg', '.jpeg', '.png', '.webp']:
                p_out = process_image_compress(u_comp, profile["quality"], profile["max_width"])
            else:
                p_in = save_uploaded_file(u_comp)
                cleanup_later.append(p_in)
                if ext == '.pdf':
                    if split_for_ai:
                        p_out, part_paths, zip_path = compress_and_split_pdf(p_in, target_total_mb, max_part_mb, pdf_mode)
                    else:
                        p_out = compress_pdf_to_target(p_in, target_total_mb, pdf_mode)
                elif ext in ['.docx', '.xlsx', '.pptx']:
                    p_out = process_office_recompress(p_in)

            if p_out:
                output_size, saved, ratio = get_compression_stats(input_size, p_out)
                st.success(f"Hoàn tất! Đã giảm {format_file_size(max(saved, 0))} ({max(ratio, 0):.1f}%).")
                r1, r2, r3 = st.columns(3)
                r1.metric("Trước", format_file_size(input_size))
                r2.metric("Sau nén", format_file_size(output_size))
                r3.metric("Giảm", f"{max(ratio, 0):.1f}%")

                base_name = os.path.splitext(u_comp.name)[0]
                with open(p_out, "rb") as f:
                    st.download_button("⬇️ Tải file đã nén", f, f"{base_name}_compressed{ext}")

                if part_paths:
                    st.markdown("### 📦 File nhỏ cho AI")
                    oversized = []
                    for idx, part in enumerate(part_paths, start=1):
                        part_size = os.path.getsize(part)
                        if part_size > mb_to_bytes(max_part_mb):
                            oversized.append(os.path.basename(part))
                        with open(part, "rb") as f:
                            st.download_button(
                                f"⬇️ Part {idx:03d} - {format_file_size(part_size)}",
                                f,
                                f"{base_name}_part_{idx:03d}.pdf",
                                key=f"dl_part_{idx}"
                            )
                    if zip_path:
                        with open(zip_path, "rb") as f:
                            st.download_button("⬇️ Tải tất cả parts (.zip)", f, f"{base_name}_ai_parts.zip")
                    if oversized:
                        st.warning("Một số part vẫn vượt giới hạn vì một trang đơn quá lớn. Hãy giảm target MB hoặc chọn nén sâu hơn.")

                st.toast("Xử lý file thành công!", icon="✅")

            cleanup_files(*cleanup_later)

# Footer
st.markdown("<br><hr><center><small>✨ Optimized by AI | Built with Python & Streamlit ✨</small></center>", unsafe_allow_html=True)